"""
Para generar un ejecutable standalone (.exe) de este script con PyInstaller:
1. Asegurate de tener PyInstaller instalado:
   pip install pyinstaller

2. Ejecutá en la terminal:
   pyinstaller --noconsole --onefile --name ComparadorPPTX pptx_comparador_gui.py

Esto creará un ejecutable en la carpeta /dist que podrás distribuir sin necesidad de instalar Python.
"""

import os
import hashlib
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
from difflib import SequenceMatcher
from collections import defaultdict
import pandas as pd

# Clase principal de la aplicación para comparar archivos PPTX
class PPTXComparadorApp:
    def __init__(self, root):
        # Inicializa la ventana principal
        self.root = root
        self.root.title("🧠 Comparador de PowerPoints")  # Título de la ventana
        self.root.geometry("1000x600")  # Tamaño de la ventana

        self.similares = []  # Lista que almacena los resultados de archivos similares
        self.archivos = []   # Lista que almacena los archivos pptx encontrados en la carpeta seleccionada
        self.carpeta_base = ""  # Carpeta de base elegida, para calcular rutas relativas

        # Botón para seleccionar la carpeta donde se encuentran los archivos PPTX
        self.btn_seleccionar = tk.Button(root, text="📂 Seleccionar Carpeta", command=self.seleccionar_carpeta)
        self.btn_seleccionar.pack(pady=10)

        # Etiqueta para mostrar información sobre la carpeta seleccionada
        self.lbl_info = tk.Label(root, text="No hay carpeta seleccionada")
        self.lbl_info.pack()

        # Tabla para mostrar los archivos comparados, con las columnas: archivo1, archivo2, similitud y borrar
        self.tree = ttk.Treeview(root, columns=("archivo1", "archivo2", "similitud", "borrar"), show="headings", selectmode="extended")
        for col in ("archivo1", "archivo2", "similitud", "borrar"):
            self.tree.heading(col, text=col)  # Establece el nombre de cada columna
            self.tree.column(col, width=250 if col != "similitud" else 80)  # Ajusta el tamaño de las columnas
        self.tree.pack(expand=True, fill="both")  # Agrega la tabla a la ventana

        # Botón para exportar los resultados a un archivo de Excel
        self.exportar_btn = tk.Button(root, text="📤 Exportar a Excel", command=self.exportar_excel)
        self.exportar_btn.pack(pady=5)

        # Botón para borrar los archivos seleccionados en la tabla
        self.borrar_btn = tk.Button(root, text="🗑️ Borrar seleccionados", command=self.borrar_seleccionados)
        self.borrar_btn.pack(pady=5)

    def seleccionar_carpeta(self):
        # Método para seleccionar la carpeta con archivos PPTX
        carpeta = filedialog.askdirectory()  # Abre el cuadro de diálogo para seleccionar una carpeta
        if carpeta:
            self.lbl_info.config(text=f"📁 Carpeta seleccionada: {carpeta}")  # Actualiza el texto de la etiqueta con la carpeta seleccionada
            self.carpeta_base = carpeta
            self.procesar_carpeta(carpeta)  # Procesa los archivos dentro de la carpeta

    def procesar_carpeta(self, carpeta):
        # Procesa todos los archivos pptx en la carpeta seleccionada
        self.similares.clear()  # Limpia la lista de archivos similares
        self.tree.delete(*self.tree.get_children())  # Elimina todas las filas de la tabla

        # 🔍 Buscar todos los archivos .pptx en la carpeta seleccionada
        self.archivos = self.buscar_pptx(carpeta)
        hashes = defaultdict(list)  # Diccionario para almacenar los archivos agrupados por hash
        contenidos = {}  # Diccionario para almacenar el texto extraído de cada archivo

        # 🧾 Recorremos todos los archivos y realizamos dos procesos:
        # 1. Calcular el hash MD5 para detectar duplicados exactos
        # 2. Extraer el texto de cada diapositiva para comparar su contenido textual
        for archivo in self.archivos:
            h = self.hash_md5(archivo)  # Calcula el hash MD5 del archivo
            hashes[h].append(archivo)  # Agrupa el archivo por su hash
            contenidos[archivo] = self.extraer_texto(archivo)  # Extrae el texto del archivo

        # 🧩 Comparación por HASH: si hay más de un archivo con el mismo hash, son duplicados exactos
        for grupo in hashes.values():
            if len(grupo) > 1:
                # Compara todos los pares posibles dentro del grupo de archivos con el mismo hash
                for i in range(len(grupo)):
                    for j in range(i+1, len(grupo)):
                        self.agregar_resultado(grupo[i], grupo[j], 1.0)  # Añade los duplicados exactos a la lista

        # 📐 Comparación por contenido textual (diapositivas)
        # Compara todos los pares de archivos diferentes
        for i in range(len(self.archivos)):
            for j in range(i+1, len(self.archivos)):
                a1, a2 = self.archivos[i], self.archivos[j]  # Selecciona dos archivos distintos
                sim = self.similitud(contenidos[a1], contenidos[a2])  # Calcula la similitud de su contenido textual
                if sim >= 0.85 and sim < 1.0:
                    self.agregar_resultado(a1, a2, sim)  # Añade los archivos con similitud parcial a la lista

    def agregar_resultado(self, archivo1, archivo2, sim):
        # Agrega un resultado de comparación a la tabla y a la lista de resultados
        rel1 = os.path.relpath(archivo1, self.carpeta_base)  # Ruta relativa
        rel2 = os.path.relpath(archivo2, self.carpeta_base)  # Ruta relativa
        # Almacena el resultado de la comparación y lo muestra en la tabla
        self.similares.append((rel1, rel2, sim))  # Guarda el par de archivos y su similitud
        # Agrega los resultados a la tabla con el valor de similitud en porcentaje
        self.tree.insert("", "end", values=(rel1, rel2, f"{sim*100:.1f}%", rel2))  # Por defecto se propone borrar archivo2

    def borrar_seleccionados(self):
        # Método para eliminar los archivos seleccionados en la tabla
        items = self.tree.selection()  # Obtiene los elementos seleccionados
        if not items:
            messagebox.showinfo("Info", "Seleccioná al menos una fila para borrar.")  # Muestra un mensaje si no se seleccionó nada
            return

        # Confirma si realmente se quieren borrar los archivos
        confirm = messagebox.askyesno("Confirmar", "¿Seguro que querés borrar los archivos seleccionados?")
        if not confirm:
            return  # Si no se confirma, no hace nada

        errores = []  # Lista para almacenar posibles errores al intentar borrar los archivos
        for item in items:
            valores = self.tree.item(item)["values"]  # Obtiene los valores de la fila seleccionada
            archivo_borrar = os.path.join(self.carpeta_base, valores[3])  # El archivo a borrar está en la cuarta columna (borrar)
            try:
                os.remove(archivo_borrar)  # Intenta eliminar el archivo
            except Exception as e:
                errores.append((archivo_borrar, str(e)))  # Si hay error, lo almacena

        if errores:
            # Si hubo errores, muestra un mensaje con los detalles
            mensaje = "⚠️ Archivos con errores al borrar:\n" + "\n".join(f"{f}: {e}" for f, e in errores)
            messagebox.showwarning("Errores", mensaje)
        else:
            messagebox.showinfo("OK", "Archivos eliminados correctamente.")  # Muestra mensaje si se borraron correctamente
        self.tree.delete(*items)  # Elimina las filas seleccionadas de la tabla

    def exportar_excel(self):
        # Método para exportar los resultados a un archivo Excel
        if not self.similares:
            messagebox.showinfo("Sin datos", "No hay resultados para exportar.")  # Muestra mensaje si no hay resultados
            return

        # Convierte los resultados a un DataFrame de Pandas
        df = pd.DataFrame([{
            "Archivo 1": a1,
            "Archivo 2": a2,
            "Similitud (%)": round(sim * 100, 1)
        } for a1, a2, sim in self.similares])

        # Abre el cuadro de diálogo para elegir la ruta y el nombre del archivo a guardar
        ruta = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel", "*.xlsx")])
        if ruta:
            df.to_excel(ruta, index=False)  # Exporta los datos a un archivo Excel
            messagebox.showinfo("Exportado", f"Archivo guardado en:\n{ruta}")  # Muestra mensaje de éxito

    # -------------------------------
    # FUNCIONES AUXILIARES
    # -------------------------------

    def buscar_pptx(self, carpeta):
        # Busca todos los archivos .pptx en la carpeta de forma recursiva
        return [os.path.join(root, f)
                for root, _, files in os.walk(carpeta)
                for f in files if f.lower().endswith(".pptx")]

    def hash_md5(self, archivo):
        # Calcula el hash MD5 (identificador único) de un archivo binario
        h = hashlib.md5()
        with open(archivo, "rb") as f:
            for bloque in iter(lambda: f.read(4096), b""):  # Lee el archivo en bloques
                h.update(bloque)  # Actualiza el hash con cada bloque
        return h.hexdigest()  # Retorna el hash MD5

    def extraer_texto(self, archivo):
        # Extrae el texto de todas las diapositivas de un archivo pptx
        try:
            prs = Presentation(archivo)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        except Exception:
            return ""

    def similitud(self, txt1, txt2):
        # Calcula la similitud entre dos textos
        return SequenceMatcher(None, txt1, txt2).ratio()

# --- EJECUCIÓN DE LA APLICACIÓN ---
if __name__ == "__main__":
    root = tk.Tk()
    app = PPTXComparadorApp(root)
    root.mainloop()
