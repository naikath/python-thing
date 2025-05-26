import os
import hashlib
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from difflib import SequenceMatcher
from collections import defaultdict
import pandas as pd

from pptx import Presentation
from docx import Document
from openpyxl import load_workbook


def custom_prompt(parent: tk.Misc, archivo1: str, archivo2: str) -> str | None:
    """Crea prompt para elegir opciones"""
    # Crea una ventana para el prompt
    prompt_window = tk.Toplevel(parent)
    prompt_window.title("¿Qué archivo querés borrar?")
    # Bloquea interaccion de las otras ventanas de tkinter
    prompt_window.grab_set()

    # Texto en la ventana con los archivos
    texto_prompt_info = tk.Label(prompt_window, text="Elegí el archivo a borrar:", font=("Segoe UI", 10, "bold"))
    texto_prompt_info.pack(padx=20, pady=(10, 5))
    texto_archivo_1 = tk.Label(prompt_window, text=f"1️⃣ {archivo1}")
    texto_archivo_1.pack(padx=20, pady=2)
    texto_archivo_2 = tk.Label(prompt_window, text=f"2️⃣ {archivo2}")
    texto_archivo_2.pack(padx=20, pady=2)

    # Resultado del prompt
    resultado: dict[str, str | None] = {"value": None}

    # Guarda el valor y cierra la ventana
    def choose(opcion_elegida: str) -> None:
        resultado["value"] = opcion_elegida
        prompt_window.destroy()

    # Crea un frame para el layout
    frame_botones = tk.Frame(prompt_window)
    frame_botones.pack(pady=10)

    # Crea botones y asigna una funcion para elegir la opcion
    boton_borrar_archivo_1 = tk.Button(frame_botones, text="🗑️ Borrar Archivo 1", command=lambda: choose("A"))
    boton_borrar_archivo_1.pack(side="left", padx=5)
    boton_borrar_archivo_2 = tk.Button(frame_botones, text="🗑️ Borrar Archivo 2", command=lambda: choose("B"))
    boton_borrar_archivo_2.pack(side="left", padx=5)
    # Crea botón para cancelar y cerrar la ventana
    boton_cancelar = tk.Button(frame_botones, text="Cancelar", command=prompt_window.destroy)
    boton_cancelar.pack(side="left", padx=5)

    # Espera a que la ventana se cierre (al elegir una opcion)
    prompt_window.wait_window()
    return resultado["value"]


class ComparadorArchivosApp:
    """Clase principal de la aplicación"""

    def __init__(self, root: tk.Tk) -> None:
        # Inicializa variables
        # Lista que almacena los resultados de archivos similares
        self.archivos_similares: list[tuple[str, str, str, float]] = []
        # Lista que almacena los archivos encontrados en la carpeta seleccionada
        self.archivos: list[str] = []
        # Carpeta de base usada como referencia para las rutas relativas
        self.carpeta_base: str = ""
        # Inicializa la ventana
        self.setup_ui(root)

    def setup_ui(self, root: tk.Tk) -> None:
        """Inicializa ui"""
        # Inicializa la ventana principal
        self.root = root
        self.root.title("🧠 Comparador de Archivos")
        self.root.geometry("1000x600")

        # Botón para seleccionar la carpeta con los archivos
        self.boton_seleccionar = tk.Button(root, text="📂 Seleccionar Carpeta", command=self.seleccionar_carpeta)
        self.boton_seleccionar.pack(pady=10)

        # Texto con información de la carpeta seleccionada
        self.texto_info = tk.Label(root, text="No hay carpeta seleccionada")
        self.texto_info.pack()

        # Tabla con columnas agrupadas por tipo de archivo
        columnas = ("tipo", "archivo1", "archivo2", "similitud")
        self.tabla = ttk.Treeview(root, columns=columnas, show="headings", selectmode="extended")
        for col in columnas:
            # Nombre y contenido del encabezado
            self.tabla.heading(col, text=col)
            # Tamaño de cada columna
            self.tabla.column(col, width=150 if col == "tipo" else 300)
        self.tabla.pack(expand=True, fill="both")

        frame_botones = tk.Frame(root)
        frame_botones.pack(pady=10)

        # Botón para exportar los resultados a un archivo de Excel
        self.boton_exportar = tk.Button(frame_botones, text="📤 Exportar a Excel", command=self.exportar_excel)
        self.boton_exportar.pack(pady=5)
        self.boton_exportar.pack(side="left", padx=5)

        # Botón para borrar los archivos seleccionados en la tabla
        self.boton_borrar = tk.Button(frame_botones, text="🗑️ Borrar seleccionados", command=self.borrar_seleccionados)
        self.boton_borrar.pack(side="left", padx=5)

    def seleccionar_carpeta(self) -> None:
        """Selecciona y procesa una carpeta si es seleccionada"""
        # Selecciona una carpeta con un explorador de archivos
        carpeta = filedialog.askdirectory()
        if carpeta:
            # Actualiza texto con la carpeta seleccionada
            self.texto_info.config(text=f"📁 Carpeta seleccionada: {carpeta}")
            self.carpeta_base = carpeta
            # Procesa los archivos dentro de la carpeta
            self.procesar_carpeta(carpeta)

    def procesar_carpeta(self, carpeta: str) -> None:
        """Procesa la carpeta con todos los metodos"""
        # Limpia la lista de archivos similares existentes
        self.archivos_similares.clear()
        # Elimina todas las filas de la tabla
        self.tabla.delete(*self.tabla.get_children())

        # 🔍 Buscar todos los archivos en la carpeta seleccionada
        self.archivos = self.buscar_archivos(carpeta)

        self.procesar_archivos()

        self.actualizar_tabla()

    def procesar_archivos(self) -> None:
        """Procesa y agrupa los archivos similares"""
        # Diccionario para almacenar los archivos agrupados por hash
        hashes: defaultdict[str, list[str]] = defaultdict(list)
        # Diccionario para almacenar el texto extraído de cada archivo
        contenidos = {}

        # 🧾 Recorremos todos los archivos y realizamos dos procesos:
        # 1. Calcular el hash MD5 para detectar duplicados exactos
        # 2. Extraer y comparar el contenido del texto
        for archivo in self.archivos:
            # Calcula el hash MD5 del archivo
            hash = self.hash_md5(archivo)
            # Agrupa el archivo por su hash
            hashes[hash].append(archivo)

            # Extrae el texto del archivo
            contenidos[archivo] = self.extraer_texto(archivo)

        # 🧩 Comparación por HASH: si hay más de un archivo con el mismo hash, son duplicados exactos
        for grupo in hashes.values():
            if len(grupo) > 1:
                # Compara todos los pares posibles dentro del grupo de archivos con el mismo hash
                for i in range(len(grupo)):
                    for j in range(i + 1, len(grupo)):
                        # Añade el primero en la lista de los hash y todos los demas como duplicados
                        self.agregar_resultado(grupo[i], grupo[j], 1.0)

        # 📐 Comparación por contenido textual (diapositivas)
        # Compara todos los pares de archivos diferentes
        for i in range(len(self.archivos)):
            for j in range(i + 1, len(self.archivos)):
                # Selecciona dos archivos distintos
                a1, a2 = self.archivos[i], self.archivos[j]

                # Compara sólo entre archivos del mismo tipo
                tipo1 = os.path.splitext(a1)[1].lower()
                tipo2 = os.path.splitext(a2)[1].lower()
                if tipo1 != tipo2:
                    continue

                # Calcula la similitud del texto
                texto1 = self.limpiar_texto(contenidos[a1])
                texto2 = self.limpiar_texto(contenidos[a2])
                similitud = self.similitud_texto(texto1, texto2)
                # Agrega el resultado en un cierto rango de coincidencia
                # No se incluye el 1.0 ya que estaría duplicado con la comparación por hash
                if 0.0 <= similitud < 1.0:
                    # Añade los archivos con similitud parcial a la lista
                    self.agregar_resultado(a1, a2, similitud)

    def agregar_resultado(self, archivo1: str, archivo2: str, similitud: float) -> None:
        """Agrega un resultado de la comparación a la lista"""
        # Rutas relativas
        rel1 = os.path.relpath(archivo1, self.carpeta_base)
        rel2 = os.path.relpath(archivo2, self.carpeta_base)
        # Tipo de archivo
        tipo: str = os.path.splitext(archivo1)[1].replace(".", "").upper()
        # Almacena en la lista el resultado en una tupla
        self.archivos_similares.append((tipo, rel1, rel2, similitud))

    def actualizar_tabla(self) -> None:
        """Imprime o actualiza la tabla"""
        # Ordenar por tipo, el primer elemento de la tupla
        self.archivos_similares.sort(key=lambda x: x[0])

        # Borra y vuelve a insertar las filas de la tabla
        self.tabla.delete(*self.tabla.get_children())
        for tipo, a1, a2, similitud in self.archivos_similares:
            self.tabla.insert("", "end", values=(tipo, a1, a2, f"{similitud * 100:.1f}%"))

    def borrar_seleccionados(self) -> None:
        """Método para eliminar los archivos seleccionados en la tabla"""
        # Obtiene los elementos seleccionados
        items = self.tabla.selection()
        if not items:
            # Muestra un mensaje si no se seleccionó nada
            messagebox.showinfo("Info", "Seleccioná al menos una fila para borrar.")
            return

        # Confirma si realmente se quieren borrar los archivos
        confirm = messagebox.askyesno("Confirmar", "¿Seguro que querés borrar los archivos seleccionados?")
        if not confirm:
            # Si no se confirma, no hace nada
            return

        # Lista para almacenar posibles errores al intentar borrar los archivos
        errores: list[tuple[str, str]] = []

        # Copia de los items porque vamos a ir modificando el tree
        for item in items:
            # Obtiene los valores de la fila seleccionada
            valores = self.tabla.item(item)["values"]
            _, archivo1_rel, archivo2_rel, _ = valores

            # Usamos el custom_prompt para elegir cuál borrar
            eleccion = custom_prompt(self.root, archivo1_rel, archivo2_rel)
            archivo_borrar_rel: str
            if eleccion == "A":
                archivo_borrar_rel = archivo1_rel
            elif eleccion == "B":
                archivo_borrar_rel = archivo2_rel
            else:
                continue  # Si se canceló, no hace nada

            archivo_borrar_abs = os.path.join(self.carpeta_base, archivo_borrar_rel)

            try:
                # Intenta eliminar el archivo
                os.remove(archivo_borrar_abs)
            except Exception as e:
                # Si hay error, lo almacena
                errores.append((archivo_borrar_rel, str(e)))
                continue

            # Eliminar de self.similares
            self.archivos_similares = [
                tupla  #
                for tupla in self.archivos_similares
                if archivo_borrar_rel not in tupla
            ]

            # Refrescar la tabla
            self.actualizar_tabla()

        if errores:
            # Si hubo errores, muestra un mensaje con los detalles
            mensaje = "⚠️ Archivos con errores al borrar:\n" + "\n".join(f"{f}: {e}" for f, e in errores)
            messagebox.showwarning("Errores", mensaje)
        else:
            # Muestra mensaje si se borraron correctamente
            messagebox.showinfo("OK", "Archivos eliminados correctamente.")

    def exportar_excel(self) -> None:
        """Método para exportar los resultados a un archivo Excel"""
        if not self.archivos_similares:
            # Muestra mensaje si no hay resultados
            messagebox.showinfo("Sin datos", "No hay resultados para exportar.")
            return

        # Convierte los resultados a un DataFrame de Pandas
        df = pd.DataFrame(
            [
                {"Tipo": tipo, "Archivo 1": a1, "Archivo 2": a2, "Similitud (%)": round(similitud * 100, 1)}  #
                for tipo, a1, a2, similitud in self.archivos_similares
            ]
        )

        # Abre el cuadro de diálogo para elegir la ruta y el nombre del archivo a guardar
        ruta = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel", "*.xlsx")])
        if ruta:
            # Exporta los datos a un archivo Excel
            df.to_excel(ruta, index=False)
            # Muestra mensaje de éxito
            messagebox.showinfo("Exportado", f"Archivo guardado en:\n{ruta}")

    # -------------------------------
    # FUNCIONES AUXILIARES
    # -------------------------------

    def buscar_archivos(self, carpeta: str) -> list[str]:
        """Busca todos los archivos en la carpeta de forma recursiva"""
        extensiones = (".pptx", ".docx", ".xlsx")
        return [
            os.path.join(root, file)  #
            for root, _, files in os.walk(carpeta)
            for file in files
            if file.lower().endswith(extensiones)
        ]

    def hash_md5(self, archivo: str) -> str:
        """Calcula el hash MD5 (identificador único) de un archivo binario"""
        hash = hashlib.md5()
        with open(archivo, "rb") as f:
            # Lee el archivo en bloques
            for bloque in iter(lambda: f.read(4096), b""):
                # Actualiza el hash con cada bloque
                hash.update(bloque)
        # Retorna el hash MD5
        return hash.hexdigest()

    def extraer_texto(self, archivo: str) -> str:
        """Detecta tipo de archivo y aplica extractor correspondiente"""
        if archivo.lower().endswith(".pptx"):
            return self.extraer_texto_pptx(archivo)
        elif archivo.lower().endswith(".docx"):
            return self.extraer_texto_docx(archivo)
        elif archivo.lower().endswith(".xlsx"):
            return self.extraer_texto_xlsx(archivo)
        return ""

    def extraer_texto_pptx(self, archivo: str) -> str:
        """Extrae el texto de todas las diapositivas de un archivo pptx"""
        prs = Presentation(archivo)
        return "\n".join(
            getattr(shape, "text")  #
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )

    def extraer_texto_docx(self, archivo: str) -> str:
        """Extrae el texto de todas las diapositivas de un archivo docx"""
        doc = Document(archivo)
        contenido: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                contenido.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    contenido.append(cell.text)
        return "\n".join(contenido)

    def extraer_texto_xlsx(self, archivo: str) -> str:
        """Extrae el texto de todas las diapositivas de un archivo xlsx"""
        wb = load_workbook(archivo, data_only=True)
        contenido: list[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                contenido.extend(
                    [
                        str(cell)  #
                        for cell in row
                        if cell is not None
                    ]
                )
        return "\n".join(contenido)

    def limpiar_texto(self, texto: str) -> str:
        """Removes any whitespace character and joins with spaces to a single string"""
        return " ".join(texto.lower().split())

    def similitud_texto(self, texto1: str, texto2: str) -> float:
        """Calcula la similitud entre dos textos"""
        return SequenceMatcher(None, texto1, texto2).ratio()


# --- EJECUCIÓN DE LA APLICACIÓN ---
if __name__ == "__main__":
    root = tk.Tk()
    app = ComparadorArchivosApp(root)
    root.mainloop()
